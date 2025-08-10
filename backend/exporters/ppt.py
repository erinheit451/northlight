from io import BytesIO
from pptx import Presentation

def build_ppt(d: dict, title: str = "Campaign Benchmark Report") -> BytesIO:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    subtitle = slide.placeholders[1]
    subtitle.text = f"CPL: {d.get('derived',{}).get('cpl')}, CPC: {d.get('derived',{}).get('cpc')}"
    bio = BytesIO()
    prs.save(bio)
    bio.seek(0)
    return bio
