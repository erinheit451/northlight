from fastapi import FastAPI, HTTPException
from fastapi.middleware.cors import CORSMiddleware
from pydantic import BaseModel
from typing import Optional, Dict, Any, List, Tuple
from pathlib import Path
import math, json
from io import BytesIO
from fastapi.responses import StreamingResponse
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from starlette.staticfiles import StaticFiles

# --- Import ALL your routers ---
from backend.routers import diagnose, export
# Import the book router
from backend.routers.book import router as book_router
from backend.data.snapshots import latest_bench_path

# --- App Setup ---
ROOT = Path(__file__).resolve().parent
from backend.data.snapshots import latest_bench_path



app = FastAPI(title="Northlight Unified API", version="0.7.0")

app.add_middleware(
    CORSMiddleware,
    allow_origins=[
        "https://northlight.pages.dev",
        "https://develop.northlight.pages.dev",
        "http://localhost",
        "http://localhost:8000",
        "http://127.0.0.1:5500",
        "http://127.0.0.1:8000",
        "http://127.0.0.1:9000",
    ],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# --- Include ALL API Routers ---
app.include_router(diagnose.router)
app.include_router(export.router)
app.include_router(book_router)  # This enables /api/book/ endpoints

# --- Static File Mounting ---
app.mount("/book", StaticFiles(directory=str(ROOT.parent / "frontend" / "book"), html=True), name="book")
app.mount("/", StaticFiles(directory=str(ROOT.parent / "frontend"), html=True), name="static")

# --- Benchmark Data Loading ---
BENCH: Dict[str, Any] = {"_version": "boot"}
TOL = 0.10

def load_benchmarks() -> Dict[str, Any]:
    """
    Loads the newest backend/data/YYYY-MM-DD-benchmarks.json and returns a dict.
    Expected shape: { "<key>": { "meta": { category, subcategory }, ... }, ... }
    """
    path = latest_bench_path()
    if not path.exists():
        raise FileNotFoundError(f"No benchmark snapshots found in {path.parent}")
    payload = json.loads(path.read_text(encoding="utf-8-sig"))
    recs = payload.get("records")
    if not isinstance(recs, dict):
        raise ValueError("Benchmark snapshot missing 'records' object")
    recs["_version"] = payload.get("version") or payload.get("date") or path.stem
    return recs
@app.on_event("startup")
def _startup():
    import logging
    global BENCH
    try:
        BENCH = load_benchmarks()
    except Exception:
        logging.exception("Failed to load benchmarks at startup; serving empty meta")
        BENCH = {"_version": "empty"}
import logging
@app.get("/benchmarks/meta")
def benchmarks_meta(limit: int = 2000):
    log = logging.getLogger("bench")
    items: List[Dict[str, Any]] = []
    for k, v in (BENCH or {}).items():
        if k == "_version":
            continue
        if not isinstance(v, dict):
            # Pinpoint bad data instead of 500
            log.error("Non-dict record in BENCH: key=%r type=%s value=%r", k, type(v).__name__, v)
            continue
        meta = v.get("meta") or {}
        items.append({
            "key": k,
            "category": meta.get("category"),
            "subcategory": meta.get("subcategory"),
        })
    items.sort(key=lambda x: ((x["category"] or ""), (x["subcategory"] or "")))
    return items[:limit]

@app.get("/health")
def health():
    return {"ok": True}
@app.get("/version")
def version():
    return {"version": app.version, "bench_version": (BENCH or {}).get("_version")}

# --- Diagnose Endpoint Models and Helpers ---
class DiagnoseIn(BaseModel):
    website: Optional[str] = None
    category: str
    subcategory: str
    budget: float
    clicks: float
    leads: float
    goal_cpl: Optional[float] = None
    impressions: Optional[float] = None
    dash_enabled: Optional[bool] = None

def kcat(category: str, subcategory: str) -> str:
    return f"{category.strip()}|{subcategory.strip()}"

def safe_div(n: Optional[float], d: Optional[float]) -> Optional[float]:
    try:
        if n is None or d in (None, 0):
            return None
        return n / d
    except ZeroDivisionError:
        return None

def r2(x): return None if x is None else round(float(x), 2)
def r4(x): return None if x is None else round(float(x), 4)

def clamp(x, lo, hi):
    return None if (x is None or lo is None or hi is None) else max(lo, min(x, hi))

# Anchors for cost metrics (CPL/CPC)
ANCHOR_MAP = [
    ("top10", 0.90),
    ("top25", 0.75),
    ("avg",   0.50),
    ("bot25", 0.25),
    ("bot10", 0.10),
]

ANCHOR_MAP_RAW = [
    ("top10", 0.10),
    ("top25", 0.25),
    ("avg",   0.50),
    ("bot25", 0.75),
    ("bot10", 0.90),
]

def _sorted_anchors_from_dms_raw(dms: Dict[str,float]) -> List[Tuple[float,float]]:
    if not dms:
        return []
    pairs=[]
    for k,p in ANCHOR_MAP_RAW:
        v=dms.get(k)
        if isinstance(v,(int,float)) and v>0:
            pairs.append((float(v), float(p)))
    pairs.sort(key=lambda t:t[0])
    return pairs

def _sorted_anchors_from_dms(dms: Dict[str,float]) -> List[Tuple[float,float]]:
    if not dms: 
        return []
    pairs=[]
    for k,p in ANCHOR_MAP:
        v=dms.get(k)
        if isinstance(v,(int,float)) and v>0: 
            pairs.append((float(v), float(p)))
    pairs.sort(key=lambda t:t[0])
    return pairs

def _percentile_from_value_log(value: float, anchors: List[Tuple[float,float]]) -> Optional[float]:
    if value is None or len(anchors) < 2:
        return None
    lo_v, lo_p = anchors[0]
    hi_v, hi_p = anchors[-1]

    if value <= lo_v:
        return round(min(1.0, lo_p + 0.05), 2)

    if value >= hi_v:
        return round(max(0.0, hi_p - 0.05), 2)

    lv = math.log(value)
    for i in range(len(anchors) - 1):
        v0, p0 = anchors[i]
        v1, p1 = anchors[i + 1]
        if v0 <= value <= v1:
            lv0, lv1 = math.log(v0), math.log(v1)
            t = 0.0 if lv1 == lv0 else (lv - lv0) / (lv1 - lv0)
            return round(p0 + (p1 - p0) * t, 2)
    return None

def value_at_percentile_log(p: float, dms: Dict[str, float]) -> Optional[float]:
    if p is None or not isinstance(p, (int, float)):
        return None
    anchors = _sorted_anchors_from_dms_raw(dms)
    if len(anchors) < 2:
        return None

    pv = [(perc, val) for (val, perc) in anchors]
    pv.sort(key=lambda t: t[0])

    lo_p, lo_v = pv[0]
    hi_p, hi_v = pv[-1]
    if p <= lo_p:
        return float(lo_v)
    if p >= hi_p:
        return float(hi_v)

    for i in range(len(pv) - 1):
        p0, v0 = pv[i]
        p1, v1 = pv[i + 1]
        if p0 <= p <= p1:
            t = 0.0 if (p1 - p0) == 0 else (p - p0) / (p1 - p0)
            lv0, lv1 = math.log(v0), math.log(v1)
            return float(math.exp(lv0 + (lv1 - lv0) * t))
    return None

def budget_anchors_from_dms(dms: Dict[str,float]) -> List[Tuple[float,float]]:
    order=[("p10_bottom",0.10),("p25_bottom",0.25),("avg",0.50),("p25_top",0.75),("p10_top",0.90)]
    pairs=[]
    for k,p in order:
        v=dms.get(k)
        if isinstance(v,(int,float)) and v>0: 
            pairs.append((float(v), float(p)))
    pairs.sort(key=lambda t:t[0])
    return pairs

def detect_edge_cases(req: DiagnoseIn, cpl: Optional[float], cpc: Optional[float], 
                      cr: Optional[float], med_cpl: Optional[float]) -> List[str]:
    issues = []
    
    if req.leads == 0 and req.clicks > 50:
        issues.append("zero_conversions_high_clicks")
    
    if cpl and med_cpl and cpl > med_cpl * 5:
        issues.append("cpl_extreme_outlier")
    
    if cr and cr < 0.001:
        issues.append("cr_extremely_low")
    
    if req.budget < 100 and req.clicks > 1000:
        issues.append("implausible_budget_clicks")
    
    return issues

@app.post("/diagnose")
def diagnose(req: DiagnoseIn):
    key = kcat(req.category, req.subcategory)
    if key not in BENCH:
        raise HTTPException(status_code=404, detail="Category/subcategory not found in benchmarks")
    bench = BENCH[key]

    # medians
    med_cpl = bench.get("cpl", {}).get("median")
    med_cpc = bench.get("cpc", {}).get("median")
    med_ctr = bench.get("ctr", {}).get("median")
    med_budget = bench.get("budget", {}).get("median")

    # DMS anchors
    cpl_dms = bench.get("cpl", {}).get("dms", {}) or {}
    cpc_dms = bench.get("cpc", {}).get("dms", {}) or {}
    budget_dms = bench.get("budget", {}).get("dms", {}) or {}

    # derived metrics
    cpc = safe_div(req.budget, req.clicks)
    cpl = safe_div(req.budget, req.leads)
    cr  = safe_div(req.leads, req.clicks)
    ctr = safe_div(req.clicks, req.impressions) if req.impressions else None

    # edge cases
    edge_cases = detect_edge_cases(req, cpl, cpc, cr, med_cpl)

    # buffers
    buffer = max(1.0, med_cpl * 0.05) if med_cpl else 0.0
    tol_hit = max(1.0, (req.goal_cpl * 0.05)) if req.goal_cpl else 1.0

    # market difficulty
    def market_band_for_goal(goal: Optional[float]) -> str:
        if goal is None:
            return "unknown"
        if med_cpl is not None and goal >= (med_cpl - buffer):
            return "acceptable"
        top25 = cpl_dms.get("top25")
        if top25 is not None and goal >= top25:
            return "aggressive"
        return "unrealistic"
    market_band = market_band_for_goal(req.goal_cpl)

    # percentiles
    anchors_cpl_disp = _sorted_anchors_from_dms(cpl_dms)
    anchors_cpl_raw  = _sorted_anchors_from_dms_raw(cpl_dms)
    anchors_cpc = _sorted_anchors_from_dms(cpc_dms)

    prob_goal = _percentile_from_value_log(req.goal_cpl, anchors_cpl_raw) if (req.goal_cpl and anchors_cpl_raw) else None

    realistic_low  = med_cpl
    realistic_high = cpl_dms.get("bot25") if cpl_dms else None

    rec66 = value_at_percentile_log(0.66, cpl_dms) if cpl_dms else None
    recommended_cpl = clamp(rec66, realistic_low, realistic_high) if rec66 is not None else (med_cpl if med_cpl is not None else None)

    cpl_pct = _percentile_from_value_log(cpl, anchors_cpl_disp) if (cpl is not None and anchors_cpl_disp) else None

    # goal status
    goal_status = "unknown"
    if req.goal_cpl is not None and cpl is not None:
        if cpl <= req.goal_cpl + tol_hit:
            goal_status = "achieved"
        elif cpl <= req.goal_cpl + 2 * tol_hit:
            goal_status = "on_track"
        else:
            goal_status = "behind"

    def band_cost_dms(value, dms, median):
        if value is None or not dms or median is None:
            return "unknown"
        bot25 = dms.get("bot25")
        if bot25 is None:
            if value <= median * (1 - TOL): return "green"
            if value >= median * (1 + TOL): return "red"
            return "amber"
        if value <= median:
            return "green"
        if value <= bot25:
            return "amber"
        return "red"

    def band_rate(value, median):
        if value is None or median is None: return "unknown"
        if value >= median * (1 + TOL): return "green"
        if value <= median * (1 - TOL): return "red"
        return "amber"

    cpc_pct = _percentile_from_value_log(cpc, anchors_cpc) if (cpc is not None and anchors_cpc) else None
    cpc_band = band_cost_dms(cpc, cpc_dms, med_cpc)
    
    cpc_band_low  = med_cpc
    cpc_band_high = cpc_dms.get("bot25")

    def band_rate_p50_p75(value, p50, p75):
        if value is None or p50 is None or p75 is None: return "unknown"
        if value >= p75: return "green"
        if value >= p50: return "amber"
        return "red"

    # CR evaluation
    cr_eval = {"value": r4(cr), "median": None, "percentile": None, "display_percentile": None, "band": "unknown", "method": None}
    market_cpc = med_cpc if isinstance(med_cpc,(int,float)) else cpc_dms.get("avg")
    cr_p50 = None
    cr_p75 = None
    if cpl_dms and isinstance(market_cpc,(int,float)) and market_cpc > 0:
        if isinstance(cpl_dms.get("avg"), (int, float)) and cpl_dms["avg"] > 0:
            cr_med = float(market_cpc) / float(cpl_dms["avg"]); cr_p50 = cr_med; cr_eval["median"] = r4(cr_med)
        if isinstance(cpl_dms.get("top25"), (int, float)) and cpl_dms["top25"] > 0:
            cr_p75 = float(market_cpc) / float(cpl_dms["top25"])
        
        cr_anchors = []
        for (k, p_cpl) in ANCHOR_MAP:
            v_cpl = cpl_dms.get(k)
            if isinstance(v_cpl, (int, float)) and v_cpl > 0:
                v_cr = float(market_cpc) / float(v_cpl)
                p_cr = 1.0 - p_cpl
                cr_anchors.append((v_cr, p_cr))
        cr_anchors.sort(key=lambda t: t[0])
        cr_pct = _percentile_from_value_log(cr, cr_anchors) if cr is not None else None
        cr_eval["percentile"] = r4(cr_pct) if cr_pct is not None else None
        cr_eval["display_percentile"] = r4(cr_pct) if cr_pct is not None else None
        cr_eval["method"] = "derived_from_cpl_dms_with_market_cpc"

    cr_band = band_rate_p50_p75(cr, cr_p50, cr_p75)
    cr_eval["band"] = cr_band
    if cr_p50 is not None and cr_p75 is not None:
        cr_eval["target_range"] = {"low": r4(cr_p50), "high": r4(cr_p75)}

    cpl_pct = _percentile_from_value_log(cpl, anchors_cpl_disp) if (cpl is not None and anchors_cpl_disp) else None
    cpl_band = band_cost_dms(cpl, cpl_dms, med_cpl)

    cpl_eval = {
        "value": r2(cpl),
        "median": r2(med_cpl),
        "percentile": r4(cpl_pct) if cpl_pct is not None else None,
        "display_percentile": r4(cpl_pct) if cpl_pct is not None else None,
        "band": cpl_band,
        "performance_tier": "strong" if cpl_band == "green" else ("average" if cpl_band == "amber" else "weak"),
        "label": "lower is better",
        "target_range": {"low": r2(realistic_low) if realistic_low is not None else None,
                         "high": r2(realistic_high) if realistic_high is not None else None}
    }
    cpc_eval = {
        "value": r2(cpc),
        "median": r2(med_cpc),
        "percentile": r4(cpc_pct) if cpc_pct is not None else None,
        "display_percentile": r4(cpc_pct) if cpc_pct is not None else None,
        "band": cpc_band,
        "performance_tier": "strong" if cpc_band == "green" else ("average" if cpc_band == "amber" else "weak"),
        "target_range": {"low": r2(cpc_band_low) if cpc_band_low is not None else None,
                         "high": r2(cpc_band_high) if cpc_band_high is not None else None}
    }

    budget_eval = {"value": r2(req.budget), "median": r2(med_budget), "percentile": None, "band": "context", "note": None}
    if isinstance(req.budget, (int, float)) and budget_dms:
        b_pct = _percentile_from_value_log(req.budget, budget_anchors_from_dms(budget_dms))
        budget_eval["percentile"] = r4(b_pct) if b_pct is not None else None
        if b_pct is not None:
            budget_eval["band"] = "above_avg" if b_pct > 0.60 else ("below_avg" if b_pct < 0.40 else "average")
    elif med_budget is not None:
        budget_eval["note"] = "Peer percentiles unavailable; showing median only."

    tags = []
    if ctr is not None and med_ctr is not None:
        if ctr < med_ctr * 0.8:
            tags.append("ad_relevance_low")
        elif cr is not None and cr < 0.05:
            tags.append("message_lp_mismatch")

    if "zero_conversions_high_clicks" in edge_cases:
        tags.append("tracking_issue_suspected")
    if "cpl_extreme_outlier" in edge_cases:
        tags.append("performance_outlier")

    def determine_goal_scenario(user_goal: Optional[float], realistic_range: Dict[str, Optional[float]], cpl_dms: Dict[str, float]) -> str:
        if user_goal is None:
            return "unknown"
        
        range_low = realistic_range.get("low")
        range_high = realistic_range.get("high")
        
        if range_low is not None and user_goal < range_low:
            return "goal_too_aggressive"
        
        p90_high = cpl_dms.get("top10")
        if p90_high is not None and user_goal > p90_high:
            return "goal_conservative"
        
        return "goal_in_range"

    goal_scenario = determine_goal_scenario(req.goal_cpl, {
        "low": realistic_low,
        "high": realistic_high
    }, cpl_dms)
    
    if goal_scenario is None:
        goal_scenario = "unknown"

    goal_analysis = {
        "market_band": market_band,
        "prob_leq_goal": r4(prob_goal) if prob_goal is not None else None,
        "recommended_cpl": r2(recommended_cpl) if recommended_cpl is not None else None,
        "realistic_range": {
            "low": r2(realistic_low) if realistic_low is not None else None,
            "high": r2(realistic_high) if realistic_high is not None else None
        },
        "goal_scenario": goal_scenario,
        "note": None,
        "can_autoadopt": market_band in ("aggressive", "unrealistic")
    }
    if prob_goal is not None and req.goal_cpl is not None:
        goal_analysis["note"] = f"{int(round((1 - prob_goal) * 100))}% of campaigns do not achieve the stated goal."

    # diagnosis logic
    primary = None
    reason = "ok"
    extra: Dict[str, Any] = {}

    goal_gap = (cpl - req.goal_cpl) if (cpl is not None and req.goal_cpl is not None) else None
    median_gap = (cpl - med_cpl) if (cpl is not None and med_cpl is not None) else None

    def choose_vs_goal():
        notes = {}
        if req.goal_cpl is None or cpl is None or cpc is None or cr is None:
            return None, notes
        cr_needed = safe_div(cpc, req.goal_cpl)
        cpc_needed = (req.goal_cpl * cr) if cr is not None else None
        notes["targets"] = {"cr_needed": cr_needed, "cpc_needed": cpc_needed}
        cpc_top10 = cpc_dms.get("top10")
        if cpc_needed is not None and cpc_top10 is not None and cpc_needed < cpc_top10:
            notes["feasibility"] = "cpc_below_top10_unrealistic"
            return "cr", notes
        rel_cr = ((cr_needed - cr) / cr) if (cr and cr_needed) else float("inf")
        rel_cpc = ((cpc - cpc_needed) / cpc) if (cpc and cpc_needed) else float("inf")
        notes["delta"] = {"rel_cr": rel_cr, "rel_cpc": rel_cpc}
        return ("cr" if rel_cr <= rel_cpc else "cpc"), notes

    celebration = None
    if goal_status == "achieved" and market_band == "aggressive":
        celebration = "exceeded_aggressive"
    elif goal_status == "achieved" and market_band == "unrealistic":
        celebration = "exceeded_unrealistic"
    elif goal_status == "achieved":
        celebration = "goal_achieved"

    cpl_gap_pct = None
    if req.goal_cpl and cpl:
        cpl_gap_pct = max(0.0, (req.goal_cpl - cpl) / req.goal_cpl)

    scale_conditions = [
        (goal_status == "achieved" and cpl_band == "green"),
        (goal_status == "achieved" and (cpl_gap_pct or 0) >= 0.20),
        (goal_status == "on_track" and cpl_band == "green" and market_band == "acceptable"),
    ]
    
    blockers = [
        (ctr is not None and med_ctr is not None and ctr < (med_ctr * 0.8)),
        (cr_band == "red"),
        ("tracking_issue_suspected" in tags),
    ]
    
    if any(scale_conditions) and not any(blockers):
        primary = "scale"
        reason = "performance_excellent"
        tags = list(set(tags + ["scale_budget"]))
    elif goal_status == "behind":
        if goal_gap is not None and goal_gap > tol_hit:
            primary, extra = choose_vs_goal()
            reason = "gap_to_goal"
        elif median_gap is not None and median_gap > 0:
            if med_cpc is not None and cpc is not None and cpc > med_cpc and (cr is None or cr >= 0.05):
                primary = "cpc"
            elif cr is not None and cr < 0.05:
                primary = "cr"
            else:
                primary = "cpc"
            reason = "gap_to_median"

    targets = {
        "cr_needed": r4(safe_div(cpc, req.goal_cpl) if (req.goal_cpl and cpc is not None) else None),
        "cpc_needed": r2((req.goal_cpl * cr) if (req.goal_cpl and cr is not None) else None),
    }
    if goal_status in ("achieved", "on_track"):
        targets = {"cr_needed": None, "cpc_needed": None}

    cpl_statement = None
    cpl_low = realistic_low
    cpl_high = realistic_high
    if cpl is not None and cpl_low is not None and cpl_high is not None and cpl_low < cpl_high:
        if cpl_low <= cpl <= cpl_high:
            cpl_statement = f"Your current CPL is ${r2(cpl)}; this sits within the typical range for this category."
        elif cpl < cpl_low:
            cpl_statement = f"Your current CPL is ${r2(cpl)}; this is better than the typical range (lower)."
        else:
            cpl_statement = f"Your current CPL is ${r2(cpl)}; this is above the typical range (higher)."

    budget_pct = budget_eval.get("percentile")
    budget_message = None
    if budget_pct is not None:
        if market_band == "unrealistic" and budget_pct < 0.25:
            budget_message = "⚠️ Ambitious goal may require higher budget vs peers."
        elif goal_status == "behind" and budget_pct < 0.50:
            budget_message = "Consider: budget may be limiting performance vs peers."
        elif goal_status == "achieved" and budget_pct < 0.25:
            budget_message = "💡 Strong performance at relatively low budget suggests scaling opportunity."
        else:
            budget_message = f"Budget context: {int(round(budget_pct*100))}% of peers spend less. Higher/lower spend isn't inherently better/worse."
    elif med_budget is not None:
        budget_message = f"Median peer budget: ${r2(med_budget)} (percentiles unavailable for this category)."

    scaling_preview = None
    if primary == "scale" and isinstance(req.budget, (int, float)) and cpl and cpl > 0:
        b0 = float(req.budget)
        steps = [0.20, 0.50]
        variants = [("flat CPL", 0.00), ("CPL +10%", 0.10), ("CPL +20%", 0.20)]
        out = []
        for s in steps:
            nb = b0 * (1.0 + s)
            proj = []
            for label, upl in variants:
                new_cpl = cpl * (1.0 + upl)
                proj.append({
                    "scenario": label, 
                    "cpl": r2(new_cpl),
                    "leads": r2(nb / new_cpl)
                })
            out.append({
                "budget_increase": f"+{int(s*100)}%", 
                "new_budget": r2(nb), 
                "lead_projections": proj
            })
        scaling_preview = {
            "scenarios": out,
            "disclaimer": "Projections assume similar targeting. Actual results may vary due to audience saturation or competitive factors."
        }

    if goal_status == "achieved" and goal_analysis.get("note"):
        goal_analysis["note"] += " Given your current CPL, the stated goal is already achieved."

    response = {
        "input": {
            "category": req.category, "subcategory": req.subcategory,
            "budget": r2(req.budget), "clicks": req.clicks, "leads": req.leads,
            "goal_cpl": r2(req.goal_cpl) if req.goal_cpl is not None else None,
            "impressions": req.impressions, "dash_enabled": bool(req.dash_enabled)
        },
        "goal_analysis": goal_analysis,
        "derived": {"cpc": r2(cpc), "cpl": r2(cpl), "cr": r4(cr), "ctr": r4(ctr)},
        "benchmarks": {
            "medians": {"cpl": r2(med_cpl), "cpc": r2(med_cpc), "ctr": r4(med_ctr), "budget": r2(med_budget)},
            "cpl_dms": cpl_dms, "cpc_dms": cpc_dms, "budget_dms": budget_dms,
            "cpl": cpl_eval, "cpc": cpc_eval, "cr": cr_eval, "budget": budget_eval
        },
        "goal_realism": {"band": market_band, "buffer": r2(buffer)},
        "diagnosis": {"primary": primary, "tags": tags, "reason": reason} | extra,
        "targets": targets,
        "overall": {
            "current_cpl_statement": cpl_statement,
            "goal_status": goal_status,
            "celebration": celebration
        },
        "advice": {
            "budget_message": budget_message,
            "scaling_preview": scaling_preview,
            "edge_cases": edge_cases
        },
        "meta": {"data_version": BENCH.get("_version"), "category_key": key},
    }
    return response

# --- PowerPoint Export Functions ---
def _add_title(prs, text: str):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    tx = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(1))
    tf = tx.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(28)
    p.font.bold = True
    return slide

def _add_kv(slide, y_in, label, value, w=5.0):
    box = slide.shapes.add_textbox(Inches(0.6), Inches(y_in), Inches(w), Inches(0.4))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = f"{label}: {value}"
    p.font.size = Pt(16)

def _bar(slide, x, y, w, h, pct, label_left, label_right):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(230, 230, 230)
    shape.line.fill.background()
    pct = 0.0 if pct is None else max(0.0, min(1.0, pct))
    fillw = pct * w
    fshape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(fillw), Inches(h))
    fshape.fill.solid()
    fshape.fill.fore_color.rgb = RGBColor(85, 170, 90)
    fshape.line.fill.background()
    l = slide.shapes.add_textbox(Inches(x), Inches(y - 0.25), Inches(w / 2), Inches(0.25))
    r = slide.shapes.add_textbox(Inches(x + w / 2), Inches(y - 0.25), Inches(w / 2), Inches(0.25))
    l.text_frame.text = label_left
    r.text_frame.text = label_right
    r.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT

def build_ppt(d: dict, title: str = "Campaign Benchmark Report") -> BytesIO:
    prs = Presentation()
    
    s1 = _add_title(prs, title)
    
    celebration = d.get("overall", {}).get("celebration")
    goal_status = d.get("overall", {}).get("goal_status", "unknown")
    market_band = d.get("goal_analysis", {}).get("market_band") or d.get("goal_realism", {}).get("band") or "unknown"
    
    if celebration == "exceeded_aggressive":
        _add_kv(s1, 1.0, "Status", "🎉 EXCEEDED AGGRESSIVE TARGET!")
    elif celebration == "exceeded_unrealistic":
        _add_kv(s1, 1.0, "Status", "🚀 EXCEEDED UNREALISTIC TARGET!")
    elif celebration == "goal_achieved":
        _add_kv(s1, 1.0, "Status", "✅ Goal Achieved")
    elif goal_status == "on_track":
        _add_kv(s1, 1.0, "Status", "📈 On Track")
    elif goal_status == "behind":
        _add_kv(s1, 1.0, "Status", "⚠️ Behind Goal")
    else:
        _add_kv(s1, 1.0, "Goal realism", str(market_band).title())
    
    der = d.get("derived", {})
    _add_kv(s1, 1.7, "CPL", f"${der.get('cpl'):.2f}" if der.get('cpl') is not None else "—")
    _add_kv(s1, 2.1, "CPC", f"${der.get('cpc'):.2f}" if der.get('cpc') is not None else "—")
    _add_kv(s1, 2.5, "CR", f"{(der.get('cr') * 100):.1f}%" if der.get('cr') is not None else "—")
    _add_kv(s1, 2.9, "CTR", f"{(der.get('ctr') * 100):.1f}%" if der.get('ctr') is not None else "—")

    s2 = _add_title(prs, "Performance vs Benchmarks")
    bm = d.get("benchmarks", {})
    cpl_disp = bm.get("cpl", {}).get("display_percentile", 0.5)
    cpc_disp = bm.get("cpc", {}).get("display_percentile", 0.5)
    cr_disp = bm.get("cr", {}).get("display_percentile", 0.5)
    
    _bar(s2, 0.6, 1.3, 9.0, 0.35, cpl_disp, "CPL: worse ←", "→ better")
    _bar(s2, 0.6, 1.9, 9.0, 0.35, cpc_disp, "CPC: worse ←", "→ better") 
    _bar(s2, 0.6, 2.5, 9.0, 0.35, cr_disp, "CR: worse ←", "→ better")

    s3 = _add_title(prs, "Diagnosis & Recommendations")
    diag = d.get("diagnosis", {}) or {}
    primary = diag.get("primary") or "—"
    reason = diag.get("reason")
    tags = ", ".join(diag.get("tags", [])) or "—"
    
    if primary == "scale":
        _add_kv(s3, 1.4, "Primary Recommendation", "🚀 SCALE BUDGET / ADD PRODUCTS")
        _add_kv(s3, 1.8, "Rationale", "Performance supports increased investment")
    elif primary == "cr":
        _add_kv(s3, 1.4, "Primary Focus", "Improve Conversion Rate")
        _add_kv(s3, 1.8, "Priority Areas", "Landing page, offer, user experience")
    elif primary == "cpc":
        _add_kv(s3, 1.4, "Primary Focus", "Reduce Cost Per Click")
        _add_kv(s3, 1.8, "Priority Areas", "Keywords, bidding, targeting")
    else:
        _add_kv(s3, 1.4, "Status", "No primary bottleneck identified")
        if reason: _add_kv(s3, 1.8, "Basis", reason)
    
    _add_kv(s3, 2.4, "Advisory tags", tags)
    
    if goal_status == "behind":
        t = d.get("targets", {}) or {}
        crn = t.get("cr_needed"); cpcn = t.get("cpc_needed")
        _add_kv(s3, 2.9, "CR needed", f"{crn*100:.1f}%" if isinstance(crn,(int,float)) else "—")
        _add_kv(s3, 3.3, "CPC needed", f"${cpcn:.2f}" if isinstance(cpcn,(int,float)) else "—")

    s4 = _add_title(prs, "Action Plan")
    recs = []
    
    if primary == "scale":
        recs = [
            "✅ Current performance validates scaling opportunity",
            "💰 Consider 20-50% budget increase with same targeting",
            "📊 Monitor CPL trends as volume increases",
            "🔍 Explore adjacent products/services for expansion"
        ]
    elif primary == "cr":
        recs = [
            "🎯 Landing page audit: headline/offer alignment with ads",
            "📝 Form optimization: reduce fields, improve flow",
            "🔒 Trust signals: testimonials, security badges, contact info",
            "📱 Mobile experience: ensure mobile-optimized conversion path"
        ]
    elif primary == "cpc":
        recs = [
            "🔍 Keyword audit: pause broad/irrelevant high-CPC terms",
            "🚫 Negative keywords: expand based on search terms report",
            "🎯 Bid strategy: test automated bidding if using manual",
            "📍 Targeting refinement: geo/daypart/device optimizations"
        ]
    else:
        recs = [
            "✅ Performance appears stable vs benchmarks",
            "🔍 Validate tracking: confirm lead attribution accuracy",
            "📊 Consider testing small optimizations to improve efficiency",
            "💭 Evaluate if goals align with business objectives"
        ]
    
    edge_cases = d.get("advice", {}).get("edge_cases", [])
    if "tracking_issue_suspected" in edge_cases:
        recs.insert(0, "🚨 PRIORITY: Verify conversion tracking setup")
    
    y = 1.4
    for r in recs:
        _add_kv(s4, y, "•", r, w=10.0)
        y += 0.4

    scaling_preview = d.get("advice", {}).get("scaling_preview")
    if scaling_preview and primary == "scale":
        s5 = _add_title(prs, "Scaling Scenarios")
        cur_leads = d.get("input", {}).get("leads")
        cur_cpl = der.get("cpl")
        _add_kv(
            s5, 1.2, "Current Performance",
            f"${cur_cpl:.2f} CPL at {cur_leads} leads" if (cur_cpl is not None and cur_leads is not None) else "—"
        )
        
        y = 1.8
        for scenario in scaling_preview.get("scenarios", []):
            budget_inc = scenario.get("budget_increase", "")
            new_budget = scenario.get("new_budget", 0)
            _add_kv(s5, y, f"Budget {budget_inc}", f"${new_budget:.0f}")
            y += 0.3
            
            for proj in scenario.get("lead_projections", []):
                scenario_name = proj.get("scenario", "")
                leads = proj.get("leads", 0)
                _add_kv(s5, y, f"  {scenario_name}", f"{leads:.0f} leads", w=8.0)
                y += 0.25
            y += 0.1
        
        disclaimer = scaling_preview.get("disclaimer", "")
        if disclaimer and y < 5.5:
            _add_kv(s5, y, "Note", disclaimer, w=10.0)

    bio = BytesIO()
    prs.save(bio)
    bio.seek(0)
    return bio

@app.post("/export/pptx")
def export_pptx(req: DiagnoseIn):
    result = diagnose(req)
    
    celebration = result.get("overall", {}).get("celebration")
    category = req.category
    subcategory = req.subcategory
    
    if celebration == "exceeded_aggressive":
        title_prefix = "🎉 CRUSHING IT"
    elif celebration == "exceeded_unrealistic": 
        title_prefix = "🚀 EXCEPTIONAL"
    elif result.get("overall", {}).get("goal_status") == "achieved":
        title_prefix = "✅ SUCCESS"
    elif result.get("diagnosis", {}).get("primary") == "scale":
        title_prefix = "📈 SCALE READY"
    else:
        title_prefix = "📊 BENCHMARK"
    
    title = f"{title_prefix} – {category} / {subcategory}"
    deck = build_ppt(result, title=title)
    
    filename = f"benchmark_{category}_{subcategory}.pptx".replace(" ", "_")
    return StreamingResponse(
        deck,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={"Content-Disposition": f'attachment; filename="{filename}"'}
    )